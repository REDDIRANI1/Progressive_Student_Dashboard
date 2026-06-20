import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    # Initialize Presentation
    prs = Presentation()
    
    # Set to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Common styles
    BG_COLOR = RGBColor(11, 15, 25)        # Dark Slate
    TEXT_WHITE = RGBColor(243, 244, 246)   # Primary white
    TEXT_GRAY = RGBColor(156, 163, 175)    # Secondary gray
    ACCENT_PURPLE = RGBColor(139, 92, 246) # Purple accent
    ACCENT_GREEN = RGBColor(16, 185, 129)  # Green success
    
    # Helper to set solid background color
    def set_slide_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    # Helper to add standard slide title
    def add_slide_header(slide, title_text):
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.833), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Arial'
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = TEXT_WHITE
        return title_box

    # Slide 1: Title Slide
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    
    # Add a glowing top bar shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_PURPLE
    shape.line.fill.background()
    
    # Title & Subtitle text box
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(11.833), Inches(3.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    # Title
    p_title = tf.paragraphs[0]
    p_title.text = "Progressive Student Dashboard"
    p_title.font.name = 'Arial'
    p_title.font.size = Pt(46)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_WHITE
    p_title.space_after = Pt(14)
    
    # Subtitle
    p_sub = tf.add_paragraph()
    p_sub.text = "End-to-End Verification & Walkthrough Report"
    p_sub.font.name = 'Arial'
    p_sub.font.size = Pt(20)
    p_sub.font.color.rgb = ACCENT_PURPLE
    p_sub.space_after = Pt(28)
    
    # Details
    p_det = tf.add_paragraph()
    p_det.text = "Date: June 20, 2026  |  Status: PASSED  |  Verified Roles: Student & Mentor"
    p_det.font.name = 'Arial'
    p_det.font.size = Pt(14)
    p_det.font.color.rgb = TEXT_GRAY

    # Slides content definition: (Title, Subtitle, BulletPoints, ImagePath)
    slides_data = [
        (
            "System Environment & Architecture",
            "Verification Setup Details",
            [
                "FastAPI Backend Server: Running on port 5001 with Uvicorn reload module.",
                "Vite + React Frontend: Running on port 5173 with TypeScript & Tailwind CSS.",
                "SQLite Local Database: Pre-populated using seed.py to ensure test repeatability.",
                "Automated Testing Sandbox: Evaluated using Chrome DevTools browser controls.",
                "Role Segmentation: Secure token routing for students and mentors."
            ],
            None # No image on this intro slide
        ),
        (
            "1. Unified Authentication Gateway",
            "Login Screen",
            [
                "Role-based segregation: Direct routes to Student or Mentor dashboards.",
                "User Validation: Real-time checks for correct email format.",
                "Backend Security: Salted bcrypt passwords matched against database.",
                "Session Management: JSON Web Tokens (JWT) saved upon successful login."
            ],
            "demo_assets/1_login_page.png"
        ),
        (
            "2. Student Dashboard Console",
            "General Progress & Metrics Summary",
            [
                "Summary Analytics: Live counters for completed lessons, time spent, and averages.",
                "Interactive Visualizations: Recharts charts displaying completion rates and trends.",
                "AI Recommendations: Custom next-steps generated based on performance metrics.",
                "Active Course Roster: Dynamic cards tracking ongoing course progress percentages."
            ],
            "demo_assets/2_student_dashboard.png"
        ),
        (
            "3. Granular Course tracking",
            "Interactive Progress Toggle",
            [
                "Curriculum Breakdown: View detailed modules and estimated reading times.",
                "Dynamic Progress Update: Toggle checkboxes to mark lessons as completed.",
                "Live API sync: Interacting with lesson toggles updates database values in real-time.",
                "Progress Recalculation: Course completion percentages adjust instantly in UI."
            ],
            "demo_assets/3_student_course_details.png"
        ),
        (
            "4. Mentor Cohort Management",
            "Cohort Dashboard Overview",
            [
                "Macro cohort view: Roster of all assigned student accounts and profiles.",
                "At-a-glance KPI analysis: Review overall averages and cumulative study times.",
                "Failing Progress Warnings: Automatic highlighting for courses requiring assistance.",
                "Actionable Drilldown: Instantly open detailed inspections of individual progress."
            ],
            "demo_assets/4_mentor_dashboard.png"
        ),
        (
            "5. Nested Progress Inspection",
            "Student Detail View (Mentor Mode)",
            [
                "Student Focus Mode: Access exact courses and active details for a selected student.",
                "Read-only Monitoring: Review lesson progress without altering student statistics.",
                "Responsive Layout: Panel overlay allows rapid swapping between cohort students.",
                "Navigation Integrity: Retains sidebar status options and workspace hierarchy."
            ],
            "demo_assets/5_mentor_student_details.png"
        ),
        (
            "6. Automated E2E Flow",
            "Recorded Interaction Flow",
            [
                "Automated Simulation: Tests end-to-end user navigation pathways in sequence.",
                "Flow Summary: Student Login -> Course Toggle -> Student Logout -> Mentor Login -> Inspector -> Mentor Logout.",
                "Playback asset: Saved in high fidelity as 'demo_assets/student_dashboard_flow.webp'.",
                "Execution Result: 100% test scenario completion with no console failures."
            ],
            "demo_assets/student_dashboard_flow.webp"
        )
    ]

    for title, subtitle, bullets, img_path in slides_data:
        slide = prs.slides.add_slide(blank_layout)
        set_slide_bg(slide)
        add_slide_header(slide, title)
        
        # Subtitle
        sub_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.2), Inches(11.833), Inches(0.4))
        sub_tf = sub_box.text_frame
        sub_tf.word_wrap = True
        sub_tf.margin_left = sub_tf.margin_top = sub_tf.margin_right = sub_tf.margin_bottom = 0
        sub_p = sub_tf.paragraphs[0]
        sub_p.text = subtitle.upper()
        sub_p.font.name = 'Arial'
        sub_p.font.size = Pt(12)
        sub_p.font.bold = True
        sub_p.font.color.rgb = ACCENT_PURPLE
        
        # Layout depends on whether an image exists
        if img_path and os.path.exists(img_path):
            # If WebP, convert to PNG first using Pillow
            if img_path.lower().endswith('.webp'):
                png_path = img_path[:-5] + "_frame.png"
                if not os.path.exists(png_path):
                    try:
                        from PIL import Image
                        with Image.open(img_path) as im:
                            # Save the first frame of the animation
                            im.seek(0)
                            im.convert('RGB').save(png_path, 'PNG')
                    except Exception as e:
                        print(f"Error converting WebP: {e}")
                if os.path.exists(png_path):
                    img_path = png_path

            # Left: Bullet list (width: 5.5 inches)
            text_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.5), Inches(5.0))
            tf = text_box.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
            
            for i, bullet in enumerate(bullets):
                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                p.text = "•  " + bullet
                p.font.name = 'Arial'
                p.font.size = Pt(15)
                p.font.color.rgb = TEXT_WHITE
                p.space_after = Pt(16)
                p.line_spacing = 1.2
            
            # Right: Image (width: 6.0 inches, aspect ratio maintained roughly)
            slide.shapes.add_picture(img_path, Inches(6.5), Inches(1.8), width=Inches(6.0))
        else:
            # Full width text (width: 11.8 inches)
            text_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.0), Inches(11.833), Inches(4.5))
            tf = text_box.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
            
            for i, bullet in enumerate(bullets):
                p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
                p.text = "✔  " + bullet
                p.font.name = 'Arial'
                p.font.size = Pt(18)
                p.font.color.rgb = TEXT_WHITE
                p.space_after = Pt(20)
                p.line_spacing = 1.3

    # Slide 9: Conclusion Slide
    slide = prs.slides.add_slide(blank_layout)
    set_slide_bg(slide)
    
    # Title
    add_slide_header(slide, "Verification Summary")
    
    # Large Badge Box
    badge_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(1.5))
    badge_tf = badge_box.text_frame
    badge_tf.word_wrap = True
    badge_tf.margin_left = badge_tf.margin_top = badge_tf.margin_right = badge_tf.margin_bottom = 0
    
    badge_p = badge_tf.paragraphs[0]
    badge_p.text = "VERIFICATION STATUS: PASSED"
    badge_p.font.name = 'Arial'
    badge_p.font.size = Pt(36)
    badge_p.font.bold = True
    badge_p.font.color.rgb = ACCENT_GREEN
    badge_p.space_after = Pt(10)
    
    desc_p = badge_tf.add_paragraph()
    desc_p.text = "All E2E flow testing components for students and mentors are functional and integrated."
    desc_p.font.name = 'Arial'
    desc_p.font.size = Pt(16)
    desc_p.font.color.rgb = TEXT_WHITE
    
    # Footer checklist
    checklist_box = slide.shapes.add_textbox(Inches(0.75), Inches(3.6), Inches(11.833), Inches(3.0))
    cl_tf = checklist_box.text_frame
    cl_tf.word_wrap = True
    cl_tf.margin_left = cl_tf.margin_top = cl_tf.margin_right = cl_tf.margin_bottom = 0
    
    highlights = [
        "Interactive toggles successfully synced with FastAPI backend and SQLAlchemy.",
        "Role segmentation checks verified on React Router client-side logic.",
        "Dynamic analytics (Recharts) matched pre-seeded database metrics exactly.",
        "A detailed HTML report with interactive elements is saved as demo.html."
    ]
    
    for i, highlight in enumerate(highlights):
        p = cl_tf.add_paragraph() if i > 0 else cl_tf.paragraphs[0]
        p.text = "✔  " + highlight
        p.font.name = 'Arial'
        p.font.size = Pt(15)
        p.font.color.rgb = TEXT_GRAY
        p.space_after = Pt(12)

    # Save presentation
    output_filename = "presentation.pptx"
    prs.save(output_filename)
    print(f"Presentation successfully saved to {output_filename}")

if __name__ == "__main__":
    create_presentation()
